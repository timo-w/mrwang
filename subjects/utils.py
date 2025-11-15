import os
import datetime
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from openai import AzureOpenAI



# Azure Open AI Details
client = AzureOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_KEY"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)

# Call Azure OpenAI
def generate_text(subject: str, topic: str, level: str, no_of_questions: str, no_of_choices: str, additional_info: str) -> str:
    system_prompt = """
        You are a helpful assistant that creates multiple-choice quiz documents.
        Do not include any formatting symbols in your response.
        Do not respond with any follow-up questions.

        Your response will follow a specific format:
        - For each question, begin with the question number.
        - Do not use bullet points, only new lines.
        - Under each question, insert a new line and then the possible answer beginning with the letter from A.
        - Between each question, add an additional new line.
        - Make option A the correct answer for each question.

        You will be making quizzes which secondary teachers in Scotland will be using. This means that:
        - Quizzes for S1-3s should contain questions which are answerable by 12-14 year olds.
        - National 4/5, Higher, and Advanced Higher quizzes should contain content which is applicable for those courses.

        The user prompt will contain the details for the quiz.

    """
    user_prompt = (
        f"Subject: {subject}"
        f"Topic: {topic}"
        f"Level: {level}"
        f"Number of questions: {no_of_questions}"
        f"Number of possible choices per question: {no_of_choices}"
        f"Additional information (if any): {additional_info}"
    )
    response = client.chat.completions.create(
        model=os.getenv("AZURE_DEPLOYMENT_NAME"),
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
    )
    return response.choices[0].message.content


# Generate Word Document
def create_quiz_doc(content: str, filename: str = "generated-quiz.docx") -> str:
    doc = Document()
    doc.add_heading(datetime.datetime.now().strftime("Generated Quiz - Created %H:%M on %B %d, %Y"), level=1)
    doc.add_paragraph(content)
    filepath = f"media/{filename}"
    os.makedirs("media", exist_ok=True)
    doc.save(filepath)
    return filepath


# For generating quiz from file
def extract_text_from_file(file_path: str) -> str:
    lower = file_path.lower()

    if lower.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if lower.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if lower.endswith(".pptx"):
        prs = Presentation(file_path)
        collected = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    collected.append(shape.text)
        return "\n".join(collected)

    return ""