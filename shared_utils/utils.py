import os
import re
from random import shuffle
from openai import AzureOpenAI
from docx import Document
from docx.shared import Pt as DocxPt
from pptx.util import Pt as PptxPt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from PyPDF2 import PdfReader


# Azure Open AI Details
client = AzureOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_KEY"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)


# Call Azure OpenAI
def generate_text(
    source_material: str,
    level: str,
    no_of_questions: str,
    no_of_choices: str,
    additional_info: str
) -> str:
    system_prompt = """
        You are a helpful assistant that creates multiple-choice quiz documents.
        Do not include any formatting symbols in your response.
        Do not respond with any follow-up questions.
        Use British English spellings, grammar, and conventions throughout.

        Your response will follow a specific format:
        - For each question, begin with the question number followed by a dot and the question text in one line.
        - Do not use bullet points, only new lines.
        - Under each question text, insert a new line and then the possible answer beginning with the letter from A.
        - Randomise the answer order, and below add "Answer: " followed by the letter of the correct answer. Below this, add a line with "Point: 1".

        You will be making quizzes which secondary teachers in Scotland will be using. This means that:
        - Quizzes for S1-3s should contain questions which are answerable by 12-14 year olds.
        - National 4/5, Higher, and Advanced Higher quizzes should contain content which is applicable for those courses.

        The user prompt will contain the details for the quiz.

    """
    user_prompt = f"""
        Source material:
        {source_material}

        Level: {level}
        Number of questions: {no_of_questions}
        Choices per question: {no_of_choices}
        Additional information: {additional_info}
    """
    response = client.chat.completions.create(
        model=os.getenv("AZURE_DEPLOYMENT_NAME"),
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
    )
    return response.choices[0].message.content


# Generate Word Document
def create_quiz_doc(content: str, title: str, filename: str = "generated-quiz.docx") -> str:
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(content)
    filepath = f"media/{filename}"
    os.makedirs("media", exist_ok=True)
    doc.save(filepath)
    return filepath


# For generating quiz from file
def extract_text_from_file(file_path: str) -> str:
    lower = file_path.lower()

    if lower.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if lower.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if lower.endswith(".pptx"):
        prs = Presentation(file_path)
        collected = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    collected.append(shape.text)
        return "\n".join(collected)

    return ""


# For generating printable quiz worksheets
def parse_quiz(text: str):
    questions = []

    blocks = re.split(r"\n\s*(?=\d+\.)", text.strip())

    for block in blocks:
        lines = [l.strip() for l in block.splitlines() if l.strip()]

        question_line = lines[0]

        options = []
        correct_letter = None

        for line in lines[1:]:
            if re.match(r"[A-Z]\.", line):
                options.append(line)

            elif line.startswith("Answer:"):
                correct_letter = line.split(":", 1)[1].strip()

        if not options or not correct_letter:
            continue  # skip malformed blocks safely

        correct_option = next(
            opt for opt in options if opt.startswith(correct_letter + ".")
        )

        questions.append({
            "question": question_line,
            "options": options,
            "correct": correct_option
        })

    return questions


# Helper functions for creating worksheets
def add_horizontal_rule(paragraph):
    p = paragraph._p
    pPr = p.get_or_add_pPr()

    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")

    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")       # thickness
    bottom.set(qn("w:space"), "8")    # spacing
    bottom.set(qn("w:color"), "auto")

    pBdr.append(bottom)
    pPr.append(pBdr)

def remove_initial_empty_paragraph(doc):
    if doc.paragraphs:
        first_para = doc.paragraphs[0]
        if not first_para.text.strip():
            p = first_para._element
            p.getparent().remove(p)


# For creating a printable quiz worksheet with answers on the last page
def create_worksheet_doc(quiz, title: str, filename="worksheet-quiz.docx"):
    doc = Document("shared_utils/templates/quiz_template.docx")
    remove_initial_empty_paragraph(doc)
    doc.add_heading(title, level=1)
    doc.add_paragraph("")

    answer_key = []

    for idx, q in enumerate(quiz, start=1):
        # Question (bold, size 12)
        q_para = doc.add_paragraph()
        q_run = q_para.add_run(q["question"])
        q_run.bold = True
        q_run.font.size = DocxPt(12)

        options = q["options"].copy()
        shuffle(options)

        new_letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
        correct_letter = None

        for i, opt in enumerate(options):
            text = opt.split(".", 1)[1].strip()
            letter = new_letters[i]

            opt_para = doc.add_paragraph()
            opt_run = opt_para.add_run(f"{letter}. {text}")

            if opt == q["correct"]:
                correct_letter = letter

        answer_key.append((idx, correct_letter))

        # Horizontal rule after each question block
        rule_para = doc.add_paragraph("")
        add_horizontal_rule(rule_para)

        # Blank line between questions
        doc.add_paragraph("")

    doc.add_page_break()
    doc.add_heading("Answers", level=1)

    for idx, letter in answer_key:
        ans_para = doc.add_paragraph()
        ans_run = ans_para.add_run(f"{idx}. {letter}")
        ans_run.font.size = Pt(12)

    filepath = f"media/{filename}"
    os.makedirs("media", exist_ok=True)
    doc.save(filepath)

    return filepath


# For creating a presentation quiz
def create_presentation_doc(quiz, title: str, filename="presentation-quiz.pptx"):
    prs = Presentation("shared_utils/templates/quiz_template.pptx")

    # --- Title slide ---
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    subtitle = slide.placeholders[1]
    subtitle.text = "Multiple Choice Quiz"

    # --- Question slides ---
    question_layout = prs.slide_layouts[1]  # Title + Content

    for idx, q in enumerate(quiz, start=1):
        slide = prs.slides.add_slide(question_layout)

        # Slide title = question
        slide.shapes.title.text = f"{q['question']}"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

       # Add options as lettered choices (A., B., C., ...)
        letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"

        for i, opt in enumerate(q["options"]):
            option_text = opt.split(".", 1)[1].strip()
            letter = letters[i]

            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = f"{letter}. {option_text}"
            p.level = 0
            p.font.size = PptxPt(24)


     # --- Answer slides ---
    answer_layout = prs.slide_layouts[1]  # reuse Title + Content

    for idx, q in enumerate(quiz, start=1):
        slide = prs.slides.add_slide(answer_layout)

        slide.shapes.title.text = f"Answer: {q['question']}"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        # q["correct"] looks like "B. Correct answer text"
        correct_letter, correct_text = q["correct"].split(".", 1)

        p = text_frame.paragraphs[0]
        p.text = f"{correct_letter}. {correct_text.strip()}"
        p.font.size = PptxPt(28)

        
    # --- Save file ---
    filepath = f"media/{filename}"
    os.makedirs("media", exist_ok=True)
    prs.save(filepath)

    return filepath


# Generate a meaningful quiz title based on the contents
def generate_quiz_title(source_material: str, level: str) -> str:
    system_prompt = """
    You generate short, clear educational quiz titles.
    Titles must be:
    - 5-10 words
    - Suitable for a Scottish secondary classroom
    - Use British English
    - Neutral and professional
    - No punctuation at the end
    - No quotation marks
    """

    user_prompt = f"""
    Create a suitable quiz title based on this material.

    Level: {level}

    Material:
    {source_material[:1500]}
    """

    response = client.chat.completions.create(
        model=os.getenv("AZURE_DEPLOYMENT_NAME"),
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.3
    )

    title = response.choices[0].message.content.strip()

    # Safety fallback
    if len(title) < 5 or len(title) > 80:
        return "AI-Generated Quiz"

    return title
